import streamlit as st
import pdfplumber
from pptx import Presentation
import docx
from openai import OpenAI
import textwrap
import hashlib
import re

# ============================================================
# CONFIG
# ============================================================

st.set_page_config(page_title="AI Study Engine", layout="wide")

# OpenAI client
client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])

# Cost control
MAX_CHUNK_SIZE = 2000
MAX_DOCUMENT_CHARS = 50000

# ============================================================
# UTIL: CLEAN TEXT
# Removes headers/footers/page noise to reduce tokens
# ============================================================

def clean_text(text):

    lines = text.split("\n")
    clean_lines = []

    for line in lines:
        line = line.strip()

        if len(line) < 3:
            continue

        if re.match(r'^\d+$', line):
            continue

        clean_lines.append(line)

    return "\n".join(clean_lines)

# ============================================================
# FILE TEXT EXTRACTION
# ============================================================

def extract_text(file):

    text = ""

    if file.name.endswith(".pdf"):

        with pdfplumber.open(file) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"

    elif file.name.endswith(".pptx"):

        prs = Presentation(file)

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

    elif file.name.endswith(".docx"):

        doc = docx.Document(file)

        for para in doc.paragraphs:
            text += para.text + "\n"

    elif file.name.endswith(".txt"):

        text = file.read().decode()

    return text

# ============================================================
# TEXT CHUNKING
# ============================================================

def chunk_text(text, size=MAX_CHUNK_SIZE):

    chunks = textwrap.wrap(
        text,
        width=size,
        break_long_words=False,
        replace_whitespace=False
    )

    return chunks

# ============================================================
# MAP STEP (CHEAP SUMMARIES)
# ============================================================

def summarize_chunk(chunk):

    prompt = f"""
Summarize the key ideas from this study material in 5 concise bullet points.

Content:
{chunk}
"""

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        temperature=0.2,
        messages=[
            {"role": "system", "content": "You summarize study material clearly."},
            {"role": "user", "content": prompt}
        ]
    )

    return response.choices[0].message.content

# ============================================================
# REDUCE STEP (BUILD STUDY MATERIAL)
# ============================================================

def build_study_material(summary_text):

    prompt = f"""
Using the following summarized material create structured study output.

Produce:

1. SUMMARY
2. CHEAT SHEET
3. STUDY NOTES
4. MEMORY TRICKS
5. FLASHCARDS (Q/A)
6. 10 PRACTICE QUESTIONS (MCQ)

Material:
{summary_text}
"""

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        temperature=0.3,
        messages=[
            {"role": "system", "content": "You create structured study guides."},
            {"role": "user", "content": prompt}
        ]
    )

    return response.choices[0].message.content

# ============================================================
# DOCUMENT PROCESSING PIPELINE
# Map → Reduce
# ============================================================

@st.cache_data(show_spinner=False)
def process_document(text):

    # Limit document size
    text = text[:MAX_DOCUMENT_CHARS]

    chunks = chunk_text(text)

    chunk_summaries = []

    for chunk in chunks:

        summary = summarize_chunk(chunk)

        chunk_summaries.append(summary)

    combined_summary = "\n".join(chunk_summaries)

    final_notes = build_study_material(combined_summary)

    return final_notes

# ============================================================
# STREAMLIT UI
# ============================================================

st.title("🗜 Crunch it")
st.markdown("---")
st.caption("Study Notes Generator")

st.markdown("""
Upload **PDF / PPTX / DOCX / TXT**

AI converts it into:

• Summary  
• Cheat Sheet  
• Study Notes  
• Memory Tricks  
• Flashcards  
• Practice Questions  
""")

uploaded_file = st.file_uploader(
    "Upload Study Document",
    type=["pdf","pptx","docx","txt"]
)

if uploaded_file:

    st.success("File uploaded successfully")

    if st.button("Generate Study Material"):

        with st.spinner("Extracting text from document..."):

            text = extract_text(uploaded_file)

        if len(text) < 100:

            st.error("Document does not contain enough readable text.")

            st.stop()

        text = clean_text(text)

        st.info(f"Processed text length: {len(text)} characters")

        with st.spinner("AI is generating study material..."):

            result = process_document(text)

        st.success("Study Material Ready")

        tab1, tab2, tab3 = st.tabs(["📖 Read", "⬇ Download", "🖨 Print"])

        with tab1:

            st.markdown(result)

        with tab2:

            st.download_button(
                label="Download TXT",
                data=result,
                file_name="ai_study_notes.txt"
            )

            st.download_button(
                label="Download Markdown",
                data=result,
                file_name="ai_study_notes.md"
            )

        with tab3:

            st.code(result)

            st.caption("Use browser Print → Save as PDF")

# ============================================================
# SIDEBAR
# ============================================================

st.sidebar.header("About")

st.sidebar.write("""
AI Study Engine converts documents into:

• Smart Notes  
• Cheat Sheets  
• Flashcards  
• Practice Questions  
""")

st.sidebar.write("Model: gpt-4o-mini")

st.sidebar.write("Optimized for low OpenAI cost")

st.sidebar.markdown("---")
st.sidebar.write("Max document size: 50k characters")
st.sidebar.write("")
st.markdown("---")
st.sidebar.write("ΛVICΛ tinker labs — 2026")
